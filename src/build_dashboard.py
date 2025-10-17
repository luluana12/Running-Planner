#!/usr/bin/env python3
"""
Game Review Sentiment Dashboard

A production-ready script that processes game review data to generate sentiment analysis,
aggregated insights, visualizations, and a PowerPoint summary.

Author: Data Engineer
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import argparse
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import re
from typing import Dict, List, Tuple


def load_or_create_data(path: str) -> pd.DataFrame:
    """
    Load game reviews from CSV or create synthetic data if file doesn't exist.
    
    Args:
        path: Path to the CSV file
        
    Returns:
        DataFrame with game review data
    """
    if os.path.exists(path):
        print(f"Loading data from {path}...")
        df = pd.read_csv(path)
        print(f"Loaded {len(df)} rows from existing CSV file.")
    else:
        print(f"File {path} not found. Creating synthetic dataset...")
        df = create_synthetic_data()
        print(f"Created {len(df)} synthetic rows.")
    
    return df


def create_synthetic_data() -> pd.DataFrame:
    """
    Create a synthetic dataset of game reviews for demonstration.
    
    Returns:
        DataFrame with synthetic game review data
    """
    np.random.seed(42)  # For reproducible results
    
    # Sample data
    games = [
        "The Elder Scrolls V: Skyrim", "Cyberpunk 2077", "The Witcher 3: Wild Hunt",
        "Grand Theft Auto V", "Red Dead Redemption 2", "Minecraft", "Fortnite",
        "Among Us", "Fall Guys", "Valorant", "League of Legends", "Counter-Strike 2",
        "Baldur's Gate 3", "Hogwarts Legacy", "Elden Ring", "God of War",
        "The Last of Us Part II", "Spider-Man", "Horizon Zero Dawn", "Ghost of Tsushima"
    ]
    
    genres = ["RPG", "Action", "Adventure", "FPS", "Battle Royale", "Strategy", "Puzzle", "Sports"]
    publishers = ["Bethesda", "CD Projekt Red", "Rockstar Games", "Mojang Studios", 
                  "Epic Games", "InnerSloth", "Mediatonic", "Riot Games", "Valve",
                  "Larian Studios", "Warner Bros", "FromSoftware", "Sony"]
    
    # Create synthetic reviews
    data = []
    for i in range(100):
        game = np.random.choice(games)
        genre = np.random.choice(genres)
        publisher = np.random.choice(publishers)
        
        # Generate rating (some missing values)
        rating = np.random.choice([1, 2, 3, 4, 5, np.nan], p=[0.1, 0.15, 0.2, 0.3, 0.2, 0.05])
        if rating is not np.nan and np.random.random() < 0.3:
            rating = rating * 2  # Some 0-10 scale
        
        # Generate review text (some missing)
        review_text = generate_synthetic_review_text(rating) if np.random.random() > 0.1 else ""
        
        data.append({
            'title': game,
            'genre': genre,
            'publisher': publisher,
            'rating': rating,
            'review_text': review_text
        })
    
    return pd.DataFrame(data)


def generate_synthetic_review_text(rating: float) -> str:
    """Generate synthetic review text based on rating."""
    if pd.isna(rating):
        return np.random.choice([
            "This game has potential but needs work.",
            "Mixed feelings about this one.",
            "Not sure what to think about this game."
        ])
    
    if rating <= 2:
        return np.random.choice([
            "This game is terrible and broken. Don't waste your time.",
            "Awful experience. Buggy, laggy, and boring.",
            "Hate this game. It's frustrating and unfinished.",
            "Bad graphics, terrible gameplay, and it crashes constantly."
        ])
    elif rating <= 3:
        return np.random.choice([
            "This game is okay but has some issues.",
            "Not great but not terrible either.",
            "Average game with some problems.",
            "It's fine but nothing special."
        ])
    elif rating <= 4:
        return np.random.choice([
            "Good game with minor issues.",
            "Fun and engaging gameplay.",
            "Solid game, would recommend.",
            "Great graphics and smooth gameplay."
        ])
    else:
        return np.random.choice([
            "Amazing game! Absolutely love it!",
            "Fantastic experience, highly recommend!",
            "Excellent game with beautiful graphics.",
            "Addictive and polished gameplay. Must play!",
            "Outstanding game that exceeded my expectations."
        ])


def standardize_columns(df: pd.DataFrame) -> pd.DataFrame:
    """
    Standardize column names to handle case-insensitive variants.
    
    Args:
        df: Input DataFrame
        
    Returns:
        DataFrame with standardized column names
    """
    # Create a mapping for common column name variants
    column_mapping = {}
    
    for col in df.columns:
        col_lower = col.lower().strip()
        
        if 'title' in col_lower or 'game' in col_lower:
            column_mapping[col] = 'title'
        elif 'genre' in col_lower or 'category' in col_lower:
            column_mapping[col] = 'genre'
        elif 'publisher' in col_lower or 'developer' in col_lower or 'studio' in col_lower:
            column_mapping[col] = 'publisher'
        elif 'rating' in col_lower or 'score' in col_lower:
            column_mapping[col] = 'rating'
        elif 'review' in col_lower or 'text' in col_lower or 'comment' in col_lower:
            column_mapping[col] = 'review_text'
    
    # Rename columns
    df = df.rename(columns=column_mapping)
    
    # Fill missing columns with NaN
    required_columns = ['title', 'genre', 'publisher', 'rating', 'review_text']
    for col in required_columns:
        if col not in df.columns:
            df[col] = np.nan
            print(f"Warning: Column '{col}' not found. Filling with NaN values.")
    
    return df


def normalize_rating(series: pd.Series) -> pd.Series:
    """
    Normalize ratings to 0-100 scale.
    
    Args:
        series: Series containing ratings
        
    Returns:
        Series with normalized ratings (0-100)
    """
    def normalize_single_rating(rating):
        if pd.isna(rating):
            return np.nan
        
        rating = float(rating)
        
        # Determine the scale and normalize
        if rating <= 5:
            return rating / 5 * 100
        elif rating <= 10:
            return rating / 10 * 100
        else:
            # Assume it's already on 0-100 scale, clamp if needed
            return max(0, min(100, rating))
    
    return series.apply(normalize_single_rating)


def text_sentiment_score(text: str) -> float:
    """
    Calculate sentiment score using a simple lexicon-based approach.
    
    Args:
        text: Review text
        
    Returns:
        Sentiment score from 0-100
    """
    if pd.isna(text) or text.strip() == "":
        return np.nan
    
    # Simple sentiment lexicon
    positive_words = {
        'amazing', 'awesome', 'fun', 'great', 'good', 'love', 'fantastic', 
        'excellent', 'engaging', 'addictive', 'polished', 'beautiful',
        'outstanding', 'brilliant', 'incredible', 'wonderful', 'perfect',
        'smooth', 'best', 'favorite', 'recommend', 'must play'
    }
    
    negative_words = {
        'bad', 'boring', 'buggy', 'broken', 'crash', 'hate', 'terrible',
        'awful', 'laggy', 'grind', 'frustrating', 'ugly', 'unfinished',
        'disappointing', 'horrible', 'worst', 'waste', 'don\'t buy',
        'avoid', 'trash', 'garbage', 'unplayable'
    }
    
    # Convert to lowercase and split into words
    words = re.findall(r'\b\w+\b', text.lower())
    
    # Count positive and negative words
    pos_count = sum(1 for word in words if word in positive_words)
    neg_count = sum(1 for word in words if word in negative_words)
    
    # Calculate score: 50 + (pos - neg) * 10, bounded to 0-100
    score = 50 + (pos_count - neg_count) * 10
    return max(0, min(100, score))


def categorize_sentiment(score: float) -> str:
    """
    Categorize sentiment score into buckets.
    
    Args:
        score: Sentiment score (0-100)
        
    Returns:
        Sentiment category string
    """
    if pd.isna(score):
        return "Unknown"
    elif score >= 67:
        return "Positive"
    elif score <= 33:
        return "Negative"
    else:
        return "Neutral"


def compute_aggregations(df: pd.DataFrame, min_reviews: int = 5) -> Dict[str, pd.DataFrame]:
    """
    Compute all required aggregations for the dashboard.
    
    Args:
        df: DataFrame with sentiment data
        min_reviews: Minimum reviews required for publisher analysis
        
    Returns:
        Dictionary containing all aggregated DataFrames
    """
    print("Computing aggregations...")
    
    # 1. Average score by genre
    avg_by_genre = df.groupby('genre')['sentiment_score'].mean().reset_index()
    avg_by_genre.columns = ['genre', 'avg_sentiment_0to100']
    avg_by_genre['avg_sentiment_0to100'] = avg_by_genre['avg_sentiment_0to100'].round(2)
    
    # 2. Sentiment distribution
    sentiment_dist = df['sentiment_bucket'].value_counts().reset_index()
    sentiment_dist.columns = ['sentiment_bucket', 'count']
    
    # Ensure all sentiment buckets are present
    all_buckets = ['Positive', 'Neutral', 'Negative', 'Unknown']
    for bucket in all_buckets:
        if bucket not in sentiment_dist['sentiment_bucket'].values:
            new_row = pd.DataFrame({'sentiment_bucket': [bucket], 'count': [0]})
            sentiment_dist = pd.concat([sentiment_dist, new_row], ignore_index=True)
    
    # 3. Top publishers by average sentiment (with minimum review threshold)
    publisher_stats = df.groupby('publisher').agg({
        'sentiment_score': ['mean', 'count']
    }).round(2)
    publisher_stats.columns = ['avg_sentiment', 'reviews']
    publisher_stats = publisher_stats[publisher_stats['reviews'] >= min_reviews].reset_index()
    publisher_stats = publisher_stats.sort_values('avg_sentiment', ascending=False)
    
    return {
        'avg_by_genre': avg_by_genre,
        'sentiment_dist': sentiment_dist,
        'top_publishers': publisher_stats
    }


def plot_avg_by_genre(df: pd.DataFrame, outpath: str) -> None:
    """Create bar chart of average sentiment by genre."""
    print(f"Creating genre sentiment chart: {outpath}")
    
    # Get data
    genre_data = df.groupby('genre')['sentiment_score'].mean().sort_values(ascending=True)
    
    # Create figure
    plt.figure(figsize=(10, 6))
    bars = plt.bar(range(len(genre_data)), genre_data.values)
    
    # Customize chart
    plt.title('Average Sentiment Score by Genre', fontsize=14, fontweight='bold')
    plt.xlabel('Genre', fontsize=12)
    plt.ylabel('Average Sentiment Score (0-100)', fontsize=12)
    plt.xticks(range(len(genre_data)), genre_data.index, rotation=45, ha='right')
    plt.ylim(0, 100)
    plt.grid(axis='y', alpha=0.3)
    
    # Add value labels on bars
    for i, (bar, value) in enumerate(zip(bars, genre_data.values)):
        plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                f'{value:.1f}', ha='center', va='bottom', fontsize=10)
    
    plt.tight_layout()
    plt.savefig(outpath, dpi=300, bbox_inches='tight')
    plt.close()


def plot_sentiment_pie(df: pd.DataFrame, outpath: str) -> None:
    """Create pie chart of sentiment distribution."""
    print(f"Creating sentiment pie chart: {outpath}")
    
    # Get data
    sentiment_counts = df['sentiment_bucket'].value_counts()
    
    # Create figure
    plt.figure(figsize=(8, 8))
    
    # Create pie chart
    wedges, texts, autotexts = plt.pie(sentiment_counts.values, 
                                      labels=sentiment_counts.index,
                                      autopct='%1.1f%%',
                                      startangle=90)
    
    # Customize chart
    plt.title('Sentiment Distribution', fontsize=14, fontweight='bold')
    
    # Make percentage text bold
    for autotext in autotexts:
        autotext.set_fontweight('bold')
        autotext.set_fontsize(10)
    
    plt.axis('equal')
    plt.tight_layout()
    plt.savefig(outpath, dpi=300, bbox_inches='tight')
    plt.close()


def plot_top_publishers(df: pd.DataFrame, outpath: str, min_reviews: int = 5) -> None:
    """Create bar chart of top publishers by average sentiment."""
    print(f"Creating top publishers chart: {outpath}")
    
    # Get data
    publisher_stats = df.groupby('publisher').agg({
        'sentiment_score': ['mean', 'count']
    })
    publisher_stats.columns = ['avg_sentiment', 'reviews']
    publisher_stats = publisher_stats[publisher_stats['reviews'] >= min_reviews]
    publisher_stats = publisher_stats.sort_values('avg_sentiment', ascending=True)
    
    if len(publisher_stats) == 0:
        print(f"Warning: No publishers found with at least {min_reviews} reviews")
        # Create empty chart
        plt.figure(figsize=(10, 6))
        plt.title('Top Publishers by Average Sentiment', fontsize=14, fontweight='bold')
        plt.text(0.5, 0.5, f'No publishers with ≥{min_reviews} reviews', 
                ha='center', va='center', transform=plt.gca().transAxes, fontsize=12)
        plt.axis('off')
        plt.tight_layout()
        plt.savefig(outpath, dpi=300, bbox_inches='tight')
        plt.close()
        return
    
    # Create figure
    plt.figure(figsize=(12, 6))
    bars = plt.bar(range(len(publisher_stats)), publisher_stats['avg_sentiment'].values)
    
    # Customize chart
    plt.title('Top Publishers by Average Sentiment', fontsize=14, fontweight='bold')
    plt.xlabel('Publisher', fontsize=12)
    plt.ylabel('Average Sentiment Score (0-100)', fontsize=12)
    plt.xticks(range(len(publisher_stats)), publisher_stats.index, rotation=45, ha='right')
    plt.ylim(0, 100)
    plt.grid(axis='y', alpha=0.3)
    
    # Add value labels on bars
    for i, (bar, value) in enumerate(zip(bars, publisher_stats['avg_sentiment'].values)):
        plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                f'{value:.1f}', ha='center', va='bottom', fontsize=10)
    
    plt.tight_layout()
    plt.savefig(outpath, dpi=300, bbox_inches='tight')
    plt.close()


def derive_key_takeaways(avg_by_genre: pd.DataFrame, top_publishers: pd.DataFrame, 
                        sentiment_dist: pd.DataFrame) -> List[str]:
    """
    Generate 2-4 key takeaways from the analysis.
    
    Args:
        avg_by_genre: Genre sentiment data
        top_publishers: Publisher sentiment data  
        sentiment_dist: Sentiment distribution data
        
    Returns:
        List of key takeaway strings
    """
    takeaways = []
    
    # 1. Overall sentiment distribution
    total_reviews = sentiment_dist['count'].sum()
    positive_pct = (sentiment_dist[sentiment_dist['sentiment_bucket'] == 'Positive']['count'].sum() / total_reviews * 100)
    negative_pct = (sentiment_dist[sentiment_dist['sentiment_bucket'] == 'Negative']['count'].sum() / total_reviews * 100)
    
    if positive_pct > negative_pct:
        sentiment_trend = f"Overall sentiment is positive with {positive_pct:.1f}% positive reviews"
    else:
        sentiment_trend = f"Overall sentiment is mixed with {negative_pct:.1f}% negative reviews"
    takeaways.append(sentiment_trend)
    
    # 2. Best performing genre
    if len(avg_by_genre) > 0:
        best_genre = avg_by_genre.loc[avg_by_genre['avg_sentiment_0to100'].idxmax()]
        best_genre_score = best_genre['avg_sentiment_0to100']
        takeaways.append(f"{best_genre['genre']} games have the highest average sentiment ({best_genre_score:.1f}/100)")
    
    # 3. Top publisher performance
    if len(top_publishers) > 0:
        top_publisher = top_publishers.iloc[0]
        top_publisher_score = top_publisher['avg_sentiment']
        top_publisher_reviews = top_publisher['reviews']
        takeaways.append(f"{top_publisher['publisher']} leads with {top_publisher_score:.1f}/100 sentiment across {top_publisher_reviews} reviews")
    
    # 4. Genre diversity or other insights
    if len(avg_by_genre) > 3:
        genre_range = avg_by_genre['avg_sentiment_0to100'].max() - avg_by_genre['avg_sentiment_0to100'].min()
        if genre_range > 20:
            takeaways.append(f"Significant sentiment variation across genres ({genre_range:.1f} point spread)")
    
    return takeaways[:4]  # Limit to 4 takeaways max


def build_pptx(charts: Dict[str, str], insights: List[str], outpath: str) -> None:
    """
    Create PowerPoint presentation with charts and insights.
    
    Args:
        charts: Dictionary mapping chart names to file paths
        insights: List of key insights
        outpath: Output PowerPoint file path
    """
    print(f"Creating PowerPoint presentation: {outpath}")
    
    # Create presentation
    prs = Presentation()
    
    # Add slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Game Review Sentiment — One-Day Dashboard"
    title_frame.paragraphs[0].font.size = Inches(0.3)
    title_frame.paragraphs[0].font.bold = True
    
    # Add key insights
    insights_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(3))
    insights_frame = insights_box.text_frame
    insights_frame.text = "Key Insights:"
    insights_frame.paragraphs[0].font.size = Inches(0.2)
    insights_frame.paragraphs[0].font.bold = True
    
    for insight in insights:
        p = insights_frame.add_paragraph()
        p.text = f"• {insight}"
        p.font.size = Inches(0.15)
        p.space_after = Inches(0.05)
    
    # Add charts (arranged in a 2x2 grid)
    chart_positions = [
        (Inches(5), Inches(1.5), Inches(4), Inches(2.5)),  # Top right
        (Inches(5), Inches(4), Inches(4), Inches(2.5)),    # Bottom right
        (Inches(0.5), Inches(4), Inches(4), Inches(2.5))   # Bottom left
    ]
    
    chart_files = [
        charts.get('genre_chart', ''),
        charts.get('publishers_chart', ''),
        charts.get('sentiment_chart', '')
    ]
    
    for i, (pos, chart_file) in enumerate(zip(chart_positions, chart_files)):
        if chart_file and os.path.exists(chart_file):
            slide.shapes.add_picture(chart_file, pos[0], pos[1], width=pos[2], height=pos[3])
    
    # Save presentation
    prs.save(outpath)


def main():
    """Main execution function."""
    parser = argparse.ArgumentParser(description='Game Review Sentiment Dashboard')
    parser.add_argument('--input', default='data/game_reviews.csv', 
                       help='Input CSV file path (default: data/game_reviews.csv)')
    parser.add_argument('--outdir', default='outputs', 
                       help='Output directory (default: outputs)')
    parser.add_argument('--min-reviews', type=int, default=5,
                       help='Minimum reviews for publisher analysis (default: 5)')
    
    args = parser.parse_args()
    
    # Create output directory
    os.makedirs(args.outdir, exist_ok=True)
    print(f"Output directory: {args.outdir}")
    
    # Load or create data
    df = load_or_create_data(args.input)
    
    # Standardize columns
    df = standardize_columns(df)
    print(f"Processing {len(df)} game reviews...")
    
    # Compute sentiment scores
    print("Computing sentiment scores...")
    
    # Normalize ratings and compute sentiment scores
    df['normalized_rating'] = normalize_rating(df['rating'])
    df['text_sentiment'] = df['review_text'].apply(text_sentiment_score)
    
    # Use rating if available, otherwise use text sentiment
    df['sentiment_score'] = df['normalized_rating'].fillna(df['text_sentiment'])
    df['sentiment_bucket'] = df['sentiment_score'].apply(categorize_sentiment)
    
    # Save cleaned data
    cleaned_path = os.path.join(args.outdir, 'game_reviews_cleaned_with_sentiment.csv')
    df.to_csv(cleaned_path, index=False)
    print(f"Saved cleaned data to: {cleaned_path}")
    
    # Compute aggregations
    aggregations = compute_aggregations(df, args.min_reviews)
    
    # Save CSV outputs
    csv_outputs = {
        'avg_score_by_genre.csv': aggregations['avg_by_genre'],
        'sentiment_distribution.csv': aggregations['sentiment_dist'],
        'top_publishers_avg_sentiment.csv': aggregations['top_publishers']
    }
    
    for filename, data in csv_outputs.items():
        output_path = os.path.join(args.outdir, filename)
        data.to_csv(output_path, index=False)
        print(f"Saved {filename}: {len(data)} rows")
    
    # Create charts
    chart_paths = {}
    
    # Genre chart
    genre_chart_path = os.path.join(args.outdir, 'avg_score_by_genre.png')
    plot_avg_by_genre(df, genre_chart_path)
    chart_paths['genre_chart'] = genre_chart_path
    
    # Sentiment pie chart
    sentiment_chart_path = os.path.join(args.outdir, 'sentiment_pie.png')
    plot_sentiment_pie(df, sentiment_chart_path)
    chart_paths['sentiment_chart'] = sentiment_chart_path
    
    # Publishers chart
    publishers_chart_path = os.path.join(args.outdir, 'top_publishers_avg_sentiment.png')
    plot_top_publishers(df, publishers_chart_path, args.min_reviews)
    chart_paths['publishers_chart'] = publishers_chart_path
    
    # Generate key insights
    insights = derive_key_takeaways(
        aggregations['avg_by_genre'],
        aggregations['top_publishers'],
        aggregations['sentiment_dist']
    )
    
    # Create PowerPoint
    pptx_path = os.path.join(args.outdir, 'Game_Review_Sentiment_Summary.pptx')
    build_pptx(chart_paths, insights, pptx_path)
    
    # Success summary
    print("\n" + "="*60)
    print("DASHBOARD GENERATION COMPLETE!")
    print("="*60)
    print(f"Processed {len(df)} game reviews")
    print(f"Output directory: {args.outdir}")
    print("\nGenerated files:")
    print(f"  Charts: {len(chart_paths)} PNG files")
    print(f"  Data: {len(csv_outputs)} CSV files")
    print(f"  Cleaned data: game_reviews_cleaned_with_sentiment.csv")
    print(f"  Presentation: Game_Review_Sentiment_Summary.pptx")
    print(f"\nKey insights:")
    for i, insight in enumerate(insights, 1):
        print(f"  {i}. {insight}")
    print("\nRun completed successfully!")


if __name__ == "__main__":
    main()
